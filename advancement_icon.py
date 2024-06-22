# -*- encoding: utf-8 -*-
"""为进度添加图标"""

import json
from pptx import Presentation as prstt
from pptx.util import Cm
from base import (
    P,
    PPT_DIR,
    IMAGE_DIR,
    sort_data,
    load_language_files,
    lang_file_list,
    language_data_all,
)

with open(P / "advancements_data.json", "r", encoding="utf-8") as f:
    data_adv = json.load(f)

sorted_data = sort_data(load_language_files(lang_file_list), "advancements")["en_us"]

prs_path = PPT_DIR / "advancements" / "output.pptx"
prs = prstt(prs_path)

for n, slide in enumerate(prs.slides):
    key, value = sorted_data[n]
    bg_path = IMAGE_DIR / "advancements" / f"{data_adv[key]["frame"]}.png"
    slide.shapes.add_picture(
        image_file=str(bg_path),
        left=Cm(25.82),
        top=Cm(7.3),
        width=Cm(5.2),
        height=Cm(5.2),
    )
    icon = data_adv[key]["icon"]
    icon_key = f"item.minecraft.{icon}"
    if icon_key in language_data_all:
        icon_path = IMAGE_DIR / "item" / f"{language_data_all[icon_key]}.png"
    else:
        icon_key = f"block.minecraft.{icon}"
        icon_path = IMAGE_DIR / "block" / f"{language_data_all[icon_key]}.png"
    slide.shapes.add_picture(
        image_file=str(icon_path),
        left=Cm(26.82),
        top=Cm(8.3),
        width=Cm(3.2),
        height=Cm(3.2),
    )

prs.save(PPT_DIR / "advancements" / "output.pptx")
