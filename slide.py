# -*- encoding: utf-8 -*-
"""自动化生成Minecraft标准译名列表所用脚本，列表载体为PowerPoint幻灯片。"""

import json
import re
from typing import Dict

from PIL import Image
from pptx import Presentation as prstt
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.presentation import Presentation

from base import (
    LANG_DIR,
    IMAGE_DIR,
    PPT_DIR,
    SLIDE_CONFIG,
    IGNORE_CATEGORIES,
    IGNORE_SUPPLEMENTS,
    sort_data,
    load_language_files,
    lang_list,
    lang_file_list,
    lang_list_table,
    LdataCo,
    LdataTuple,
)


def update_language_data(data: LdataCo, new_string: bool = False) -> LdataCo:
    """
    修正语言文件数据。

    Args:
        data (LdataCo): 原始语言数据
        new_string (bool): 是否为新增字符串数据

    Returns:
        LdataCo: 修正后的语言数据
    """

    updated_data = data.copy()
    for lang in lang_list:
        smithing_template_str = data[lang].get("item.minecraft.smithing_template", "")
        trim_keys = [key for key in data[lang] if "trim_smithing_template" in key]
        for key in trim_keys:
            variant = key.split(".")[2].split("_", maxsplit=1)[0]
            trim_pattern_str = data[lang].get(f"trim_pattern.minecraft.{variant}", "")
            updated_data[lang][key] = (
                f"{trim_pattern_str} {smithing_template_str}"
                if lang == "en_us"
                else f"{trim_pattern_str}{smithing_template_str}"
            )

        keys_to_delete = [
            key
            for key in data[lang]
            if key.startswith("item.minecraft.music_disc")
            or re.match(r"^item\.minecraft\.[^.]*_banner_pattern$", key)
        ]
        for key in keys_to_delete:
            updated_data[lang].pop(key, None)
        updated_data[lang].pop("item.minecraft.smithing_template", None)

        if not new_string:
            netherite_upgrade_str = data[lang].get(
                "upgrade.minecraft.netherite_upgrade", ""
            )
            banner_pattern_str = data[lang].get("item.minecraft.mojang_banner_pattern", "")
            music_disc_str = data[lang].get("item.minecraft.music_disc_5", "")
            updated_data[lang]["item.minecraft.netherite_upgrade_smithing_template"] = (
                f"{netherite_upgrade_str} {smithing_template_str}"
                if lang == "en_us"
                else f"{netherite_upgrade_str}{smithing_template_str}"
            )
            updated_data[lang]["item.minecraft.music_disc_*"] = music_disc_str
            updated_data[lang]["item.minecraft.*_banner_pattern"] = banner_pattern_str

    return updated_data


def load_supplements(data: LdataCo) -> None:
    """
    读取并合并补充字符串。

    Args:
        data (LdataCo): 语言数据
    """

    with open(LANG_DIR / "supplements.json", "r", encoding="utf-8") as f:
        supplements = json.load(f)
    for lang in lang_list_table:
        data[lang].update(supplements[lang])
    print(f"已补充{len(supplements['zh_cn'])}条字符串。")


def edit_text(
    slide_data: Dict[str, LdataTuple], category: str, prs: Presentation
) -> None:
    """
    编辑幻灯片文本。

    Args:
        slide_data (Dict[str, LdataTuple]): 幻灯片数据
        category (str): 分类名称
        prs (Presentation): 幻灯片对象
    """

    print(f"开始编辑幻灯片文本，分类：{category}。")
    for n, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                # 源字符串
                if tf.text == "Source String":
                    tf.text = slide_data["en_us"][n][1]
                    paragraph = tf.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.font.name = SLIDE_CONFIG["font"]["source"]
                    paragraph.font.size = Pt(SLIDE_CONFIG["size"]["source"])
                    paragraph.font.bold = SLIDE_CONFIG["bold"]["source"]
                # 本地化键名
                elif tf.text == "Translation Key":
                    tf.text = slide_data["en_us"][n][0]
                    paragraph = tf.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.font.name = SLIDE_CONFIG["font"]["key"]
                    paragraph.font.size = Pt(SLIDE_CONFIG["size"]["key"])
                    paragraph.font.bold = SLIDE_CONFIG["bold"]["key"]
            # 编辑表格
            if shape.has_table:
                table = shape.table
                for i, lang_name in enumerate(
                    ["zh_cn", "zh_hk", "zh_tw", "lzh"], start=1
                ):
                    cell = table.cell(i, 1)
                    cell.text = slide_data[lang_name][n][1]
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.vertical_anchor = MSO_ANCHOR.MIDDLE
                    paragraph.font.name = SLIDE_CONFIG["font"][lang_name]
                    paragraph.font.size = Pt(SLIDE_CONFIG["size"][lang_name])
                    paragraph.font.bold = SLIDE_CONFIG["bold"][lang_name]


def add_image(
    slide_data: Dict[str, LdataTuple], category: str, prs: Presentation
) -> None:
    """
    在幻灯片中添加图片。

    Args:
        slide_data (Dict[str, LdataTuple]): 幻灯片数据
        category (str): 分类名称
        prs (Presentation): 幻灯片对象
    """

    print(f"开始添加图片，分类：{category}。")
    for n, slide in enumerate(prs.slides):
        img_path = IMAGE_DIR / category / f"{slide_data['en_us'][n][1]}.png"
        if img_path.exists():
            img = Image.open(img_path)
            img_width, img_height = img.size
            img_width_cm = img_width / 72 * 2.54
            img_height_cm = img_height / 72 * 2.54

            if img_height_cm > 12:
                img_width_cm *= 12 / img_height_cm
                img_height_cm = 12
            if img_width_cm > 8:
                img_height_cm *= 8 / img_width_cm
                img_width_cm = 8

            slide.shapes.add_picture(
                image_file=str(img_path),
                left=Cm(28.57 - img_width_cm / 2),
                top=Cm((19.05 - img_height_cm) / 2 + 0.38),
                width=Cm(img_width_cm),
                height=Cm(img_height_cm),
            )


def edit_slide(slide_data: Dict[str, LdataTuple], category: str) -> None:
    """
    编辑幻灯片。

    Args:
        slide_data (Dict[str, LdataTuple]): 幻灯片数据
        category (str): 分类名称
    """
    prs_path = PPT_DIR / category / "copied.pptx"
    if prs_path.exists():
        prs = prstt(prs_path)
        edit_text(slide_data, category, prs)
        if category not in {"advancements", "biome", "enchantment"}:
            add_image(slide_data, category, prs)
        prs.save(PPT_DIR / category / "output.pptx")
    else:
        print(f"分类{category}不存在已复制的幻灯片。")


def main() -> None:
    """
    主函数，生成幻灯片内容。
    """

    data, data_all = load_language_files(lang_file_list)
    data = update_language_data(data, True)
    data_all = update_language_data(data_all)
    for lang in lang_list:
        trims = {k: data_all[lang][k] for k in data[lang] if "trim_smithing_template" in k}
        data[lang].update(trims)
    if not IGNORE_SUPPLEMENTS:
        load_supplements(data)

    print("开始生成幻灯片内容。")
    if not IGNORE_CATEGORIES["advancements"]:
        sorted_data_advancements = sort_data(data, "advancements")
        edit_slide(sorted_data_advancements, "advancements")
    if not IGNORE_CATEGORIES["biome"]:
        sorted_data_biome = sort_data(data, "biome")
        edit_slide(sorted_data_biome, "biome")
    if not IGNORE_CATEGORIES["block"]:
        sorted_data_block = sort_data(data, "block")
        edit_slide(sorted_data_block, "block")
    if not IGNORE_CATEGORIES["entity"]:
        sorted_data_entity = sort_data(data, "entity")
        edit_slide(sorted_data_entity, "entity")
    if not IGNORE_CATEGORIES["item"]:
        sorted_data_item = sort_data(data, "item", "filled_map")
        edit_slide(sorted_data_item, "item")
    if not IGNORE_CATEGORIES["effect"]:
        sorted_data_effect = sort_data(data, "effect")
        edit_slide(sorted_data_effect, "effect")
    if not IGNORE_CATEGORIES["enchantment"]:
        sorted_data_enchantment = sort_data(data, "enchantment")
        edit_slide(sorted_data_enchantment, "enchantment")
    print("已完成。")


if __name__ == "__main__":
    main()
